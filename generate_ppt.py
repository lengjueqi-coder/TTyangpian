#!/usr/bin/env python3
"""
摄影师作品集PPT生成脚本
生成精美的16:9宽屏PPT，包含封面、分类展示页和结尾页
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import os
import math

# ============ 配置 ============
BASE_DIR = '/Users/Apple/Downloads/摄影师作品/作品集'
OUTPUT_PATH = '/Users/Apple/Downloads/摄影师作品集.pptx'

# PPT尺寸 16:9 宽屏
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# 颜色方案 - 高级暗色调
BG_BLACK = RGBColor(0x0A, 0x0A, 0x0A)          # 深黑背景
BG_DARK = RGBColor(0x14, 0x14, 0x14)            # 深灰背景
ACCENT_GOLD = RGBColor(0xC8, 0xA2, 0x5C)        # 金色点缀
ACCENT_GOLD_LIGHT = RGBColor(0xD4, 0xB8, 0x7A)  # 浅金
TEXT_WHITE = RGBColor(0xF0, 0xF0, 0xF0)          # 白色文字
TEXT_LIGHT = RGBColor(0xCC, 0xCC, 0xCC)          # 浅灰文字
TEXT_DIM = RGBColor(0x88, 0x88, 0x88)            # 暗灰文字
DIVIDER_COLOR = RGBColor(0x3A, 0x3A, 0x3A)       # 分割线

# 分类配置 - 中文名 + 英文副标题 + 排列方式
CATEGORIES = [
    {
        'folder': '肖像',
        'title': '肖像',
        'subtitle': 'PORTRAIT',
        'layout': 'grid',  # 网格布局
        'max_per_slide': 6,
    },
    {
        'folder': '古风作品',
        'title': '古风作品',
        'subtitle': 'ANCIENT STYLE',
        'layout': 'grid',
        'max_per_slide': 6,
    },
    {
        'folder': '男生作品',
        'title': '男生作品',
        'subtitle': 'MENSWEAR',
        'layout': 'grid',
        'max_per_slide': 6,
    },
    {
        'folder': '艺人1',
        'title': '艺人合作',
        'subtitle': 'CELEBRITY I',
        'layout': 'featured',  # 突出展示，每张独立或2张
        'max_per_slide': 2,
    },
    {
        'folder': '艺人2',
        'title': '艺人合作',
        'subtitle': 'CELEBRITY II',
        'layout': 'featured',
        'max_per_slide': 2,
    },
    {
        'folder': '艺人3',
        'title': '艺人合作',
        'subtitle': 'CELEBRITY III',
        'layout': 'featured',
        'max_per_slide': 2,
    },
    {
        'folder': '艺人4',
        'title': '艺人合作',
        'subtitle': 'CELEBRITY IV',
        'layout': 'featured',
        'max_per_slide': 2,
    },
    {
        'folder': '品牌合作作品',
        'title': '品牌合作',
        'subtitle': 'BRAND COLLABORATION',
        'layout': 'featured',
        'max_per_slide': 1,
    },
    {
        'folder': '影视合作',
        'title': '影视合作',
        'subtitle': 'FILM & TV',
        'layout': 'featured',
        'max_per_slide': 1,
    },
]


def set_slide_bg(slide, color):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='微软雅黑'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_line(slide, left, top, width, color=ACCENT_GOLD, height=Pt(1)):
    """添加装饰线"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_image_safe(slide, img_path, left, top, width, height):
    """安全添加图片，处理RGBA等问题"""
    try:
        # 如果是PNG带alpha通道，先转换为RGB
        if img_path.lower().endswith('.png'):
            img = Image.open(img_path)
            if img.mode == 'RGBA':
                # 创建白色背景替换透明
                bg = Image.new('RGB', img.size, (10, 10, 10))
                bg.paste(img, mask=img.split()[3])
                temp_path = img_path + '_rgb.jpg'
                bg.save(temp_path, 'JPEG', quality=95)
                pic = slide.shapes.add_picture(temp_path, left, top, width, height)
                os.remove(temp_path)
                return pic
        pic = slide.shapes.add_picture(img_path, left, top, width, height)
        return pic
    except Exception as e:
        print(f"  [警告] 无法添加图片 {img_path}: {e}")
        return None


def calc_image_size(img_path, max_w, max_h):
    """计算图片在最大尺寸内的等比缩放尺寸"""
    img = Image.open(img_path)
    w, h = img.size
    ratio_w = max_w / w
    ratio_h = max_h / h
    ratio = min(ratio_w, ratio_h)
    return int(w * ratio), int(h * ratio)


def get_image_files(folder_path):
    """获取文件夹中的所有图片文件（排除.DS_Store等）"""
    files = []
    for f in sorted(os.listdir(folder_path)):
        if f.startswith('.') or f == 'Thumbs.db':
            continue
        if f.lower().endswith(('.jpg', '.jpeg', '.png', '.webp', '.bmp')):
            files.append(f)
    return files


def create_cover_slide(prs):
    """创建封面页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    set_slide_bg(slide, BG_BLACK)

    # 摄影师形象照 - 居中偏左展示
    photo_path = os.path.join(BASE_DIR, '主页摄影师形象照.jpeg')
    if os.path.exists(photo_path):
        img = Image.open(photo_path)
        # 竖版照片，放在左侧，占满高度
        max_h = Inches(6.8)
        max_w = Inches(4.5)
        new_w, new_h = calc_image_size(photo_path, max_w, max_h)
        # 居左偏上
        left = Inches(0.8)
        top = Inches(0.35)
        add_image_safe(slide, photo_path, left, top, Emu(new_w), Emu(new_h))

    # 右侧文字区域
    text_left = Inches(6.2)

    # 金色装饰线 - 上方
    add_line(slide, text_left, Inches(1.5), Inches(3.5), ACCENT_GOLD, Pt(2))

    # 主标题
    add_text_box(slide, text_left, Inches(1.8), Inches(5.5), Inches(1.2),
                 '摄影师作品集', font_size=48, color=TEXT_WHITE, bold=True,
                 alignment=PP_ALIGN.LEFT, font_name='微软雅黑')

    # 英文副标题
    add_text_box(slide, text_left, Inches(3.0), Inches(5.5), Inches(0.6),
                 'PHOTOGRAPHER PORTFOLIO', font_size=20, color=ACCENT_GOLD,
                 bold=False, alignment=PP_ALIGN.LEFT, font_name='Arial')

    # 金色装饰线 - 中间
    add_line(slide, text_left, Inches(3.8), Inches(1.5), ACCENT_GOLD, Pt(1))

    # 描述文字
    add_text_box(slide, text_left, Inches(4.1), Inches(5.5), Inches(1.5),
                 '肖像 · 古风 · 艺人 · 品牌 · 影视',
                 font_size=16, color=TEXT_LIGHT, bold=False,
                 alignment=PP_ALIGN.LEFT, font_name='微软雅黑')

    # 底部金色装饰线
    add_line(slide, Inches(0), Inches(7.2), SLIDE_WIDTH, ACCENT_GOLD, Pt(2))


def create_category_title_slide(prs, cat_info, slide_num):
    """创建分类标题页（过渡页）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # 大号中文标题 - 居中
    add_text_box(slide, Inches(0), Inches(2.0), SLIDE_WIDTH, Inches(1.5),
                 cat_info['title'], font_size=54, color=TEXT_WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name='微软雅黑')

    # 英文副标题 - 居中
    add_text_box(slide, Inches(0), Inches(3.5), SLIDE_WIDTH, Inches(0.8),
                 cat_info['subtitle'], font_size=22, color=ACCENT_GOLD,
                 bold=False, alignment=PP_ALIGN.CENTER, font_name='Arial')

    # 金色装饰线 - 上
    add_line(slide, Inches(5.5), Inches(1.8), Inches(2.333), ACCENT_GOLD, Pt(2))
    # 金色装饰线 - 下
    add_line(slide, Inches(5.5), Inches(4.5), Inches(2.333), ACCENT_GOLD, Pt(2))

    # 页码
    add_text_box(slide, Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.5),
                 f'{slide_num}', font_size=12, color=TEXT_DIM,
                 alignment=PP_ALIGN.RIGHT, font_name='Arial')


def create_grid_slide(prs, cat_info, images, slide_num):
    """创建网格布局幻灯片 - 适合肖像/古风/男生等大量图片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_BLACK)

    # 左侧标题区
    title_area_w = Inches(2.0)
    add_text_box(slide, Inches(0.3), Inches(0.5), title_area_w, Inches(0.6),
                 cat_info['title'], font_size=28, color=TEXT_WHITE, bold=True,
                 alignment=PP_ALIGN.LEFT, font_name='微软雅黑')
    add_text_box(slide, Inches(0.3), Inches(1.1), title_area_w, Inches(0.4),
                 cat_info['subtitle'], font_size=12, color=ACCENT_GOLD,
                 alignment=PP_ALIGN.LEFT, font_name='Arial')
    add_line(slide, Inches(0.3), Inches(1.6), Inches(1.5), ACCENT_GOLD, Pt(1))

    # 图片区域 - 右侧
    img_area_left = Inches(2.5)
    img_area_top = Inches(0.3)
    img_area_w = Inches(10.5)
    img_area_h = Inches(6.9)

    num_images = len(images)
    # 动态计算网格
    if num_images <= 3:
        cols = num_images
        rows = 1
    elif num_images <= 4:
        cols = 2
        rows = 2
    elif num_images <= 6:
        cols = 3
        rows = 2
    elif num_images <= 8:
        cols = 4
        rows = 2
    elif num_images <= 9:
        cols = 3
        rows = 3
    else:
        cols = 4
        rows = 3

    # 重新限制实际展示数量
    actual_count = min(num_images, cols * rows)

    cell_w = img_area_w / cols
    cell_h = img_area_h / rows
    # 图片间距
    gap = Inches(0.15)

    for i in range(actual_count):
        img_path = images[i]
        col = i % cols
        row = i // cols

        cell_left = img_area_left + Emu(int(cell_w * col))
        cell_top = img_area_top + Emu(int(cell_h * row))

        # 计算图片在cell内的等比缩放
        max_w = int(cell_w - gap * 2)
        max_h = int(cell_h - gap * 2)

        try:
            new_w, new_h = calc_image_size(img_path, max_w, max_h)
            # 居中放置
            offset_x = int((int(cell_w) - new_w) / 2)
            offset_y = int((int(cell_h) - new_h) / 2)

            pic_left = Emu(int(cell_left) + offset_x)
            pic_top = Emu(int(cell_top) + offset_y)
            add_image_safe(slide, img_path, pic_left, pic_top, Emu(new_w), Emu(new_h))
        except Exception as e:
            print(f"  [警告] 处理图片失败 {img_path}: {e}")

    # 页码
    add_text_box(slide, Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.5),
                 f'{slide_num}', font_size=12, color=TEXT_DIM,
                 alignment=PP_ALIGN.RIGHT, font_name='Arial')


def create_featured_slide(prs, cat_info, images, names, slide_num):
    """创建突出展示幻灯片 - 适合艺人/品牌等少量精选图片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_BLACK)

    num_images = len(images)

    if num_images == 1:
        # 单张 - 大幅展示
        # 左侧标题
        add_text_box(slide, Inches(0.5), Inches(0.5), Inches(3), Inches(0.6),
                     cat_info['title'], font_size=28, color=TEXT_WHITE, bold=True,
                     alignment=PP_ALIGN.LEFT, font_name='微软雅黑')
        add_text_box(slide, Inches(0.5), Inches(1.1), Inches(3), Inches(0.4),
                     cat_info['subtitle'], font_size=12, color=ACCENT_GOLD,
                     alignment=PP_ALIGN.LEFT, font_name='Arial')
        add_line(slide, Inches(0.5), Inches(1.6), Inches(1.5), ACCENT_GOLD, Pt(1))

        # 图片名
        if names and names[0]:
            display_name = names[0].replace('艺人-', '').replace('品牌合作', '')
            add_text_box(slide, Inches(0.5), Inches(2.0), Inches(3), Inches(0.5),
                         display_name, font_size=18, color=TEXT_LIGHT, bold=False,
                         alignment=PP_ALIGN.LEFT, font_name='微软雅黑')

        # 图片 - 居中大展示
        max_w = Inches(8)
        max_h = Inches(6.5)
        new_w, new_h = calc_image_size(images[0], max_w, max_h)
        pic_left = Emu(int((int(SLIDE_WIDTH) - new_w) / 2) + int(Inches(1)))
        pic_top = Emu(int((int(SLIDE_HEIGHT) - new_h) / 2))
        add_image_safe(slide, images[0], pic_left, pic_top, Emu(new_w), Emu(new_h))

    elif num_images == 2:
        # 两张 - 左右对称展示
        # 顶部标题区
        add_text_box(slide, Inches(0), Inches(0.2), SLIDE_WIDTH, Inches(0.6),
                     cat_info['title'], font_size=24, color=TEXT_WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name='微软雅黑')
        add_text_box(slide, Inches(0), Inches(0.7), SLIDE_WIDTH, Inches(0.4),
                     cat_info['subtitle'], font_size=11, color=ACCENT_GOLD,
                     alignment=PP_ALIGN.CENTER, font_name='Arial')
        add_line(slide, Inches(5.5), Inches(1.15), Inches(2.333), ACCENT_GOLD, Pt(1))

        # 左右两张图片
        for i, (img_path, name) in enumerate(zip(images, names)):
            section_w = Inches(6.2)
            section_left = Inches(0.3) + Emu(int(section_w * i))

            max_w = Inches(5.5)
            max_h = Inches(5.5)
            new_w, new_h = calc_image_size(img_path, max_w, max_h)

            offset_x = int((int(section_w) - new_w) / 2)
            pic_left = Emu(int(section_left) + offset_x)
            pic_top = Inches(1.4)

            add_image_safe(slide, img_path, pic_left, pic_top, Emu(new_w), Emu(new_h))

            # 图片名
            if name:
                display_name = name.replace('艺人-', '').replace('品牌合作', '')
                add_text_box(slide, section_left, Inches(6.8), section_w, Inches(0.5),
                             display_name, font_size=14, color=TEXT_LIGHT,
                             alignment=PP_ALIGN.CENTER, font_name='微软雅黑')

    # 页码
    add_text_box(slide, Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.5),
                 f'{slide_num}', font_size=12, color=TEXT_DIM,
                 alignment=PP_ALIGN.RIGHT, font_name='Arial')


def create_ending_slide(prs, total_slides):
    """创建结尾页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_BLACK)

    # 金色装饰线 - 上
    add_line(slide, Inches(5.5), Inches(2.5), Inches(2.333), ACCENT_GOLD, Pt(2))

    # 感谢文字
    add_text_box(slide, Inches(0), Inches(2.8), SLIDE_WIDTH, Inches(1.2),
                 'THANK YOU', font_size=48, color=TEXT_WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name='Arial')

    add_text_box(slide, Inches(0), Inches(4.0), SLIDE_WIDTH, Inches(0.8),
                 '感谢观看', font_size=24, color=ACCENT_GOLD, bold=False,
                 alignment=PP_ALIGN.CENTER, font_name='微软雅黑')

    # 金色装饰线 - 下
    add_line(slide, Inches(5.5), Inches(5.0), Inches(2.333), ACCENT_GOLD, Pt(2))

    # 底部装饰线
    add_line(slide, Inches(0), Inches(7.2), SLIDE_WIDTH, ACCENT_GOLD, Pt(2))


def main():
    print("=" * 50)
    print("摄影师作品集PPT生成器")
    print("=" * 50)

    # 创建PPT
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_num = 1

    # 1. 封面页
    print("\n[1] 创建封面页...")
    create_cover_slide(prs)
    slide_num += 1

    # 2. 各分类展示
    for cat_idx, cat_info in enumerate(CATEGORIES):
        folder_path = os.path.join(BASE_DIR, cat_info['folder'])
        if not os.path.exists(folder_path):
            print(f"\n[跳过] 文件夹不存在: {folder_path}")
            continue

        files = get_image_files(folder_path)
        if not files:
            print(f"\n[跳过] 无图片: {folder_path}")
            continue

        print(f"\n[{cat_idx+2}] 处理分类: {cat_info['title']} ({len(files)}张图片)")

        # 分类标题过渡页
        create_category_title_slide(prs, cat_info, slide_num)
        slide_num += 1

        # 根据布局类型生成内容页
        if cat_info['layout'] == 'grid':
            # 网格布局 - 每页最多展示max_per_slide张
            max_per = cat_info['max_per_slide']
            # 对于大量图片，分页展示
            chunks = []
            for i in range(0, len(files), max_per):
                chunk = files[i:i+max_per]
                chunks.append(chunk)

            for chunk_idx, chunk in enumerate(chunks):
                img_paths = [os.path.join(folder_path, f) for f in chunk]
                create_grid_slide(prs, cat_info, img_paths, slide_num)
                slide_num += 1
                print(f"  网格页 {chunk_idx+1}/{len(chunks)}: {len(chunk)}张")

        elif cat_info['layout'] == 'featured':
            # 突出展示 - 每页1-2张
            max_per = cat_info['max_per_slide']
            for i in range(0, len(files), max_per):
                chunk = files[i:i+max_per]
                img_paths = [os.path.join(folder_path, f) for f in chunk]
                names = [os.path.splitext(f)[0] for f in chunk]
                create_featured_slide(prs, cat_info, img_paths, names, slide_num)
                slide_num += 1
                print(f"  突出展示页: {names}")

    # 3. 结尾页
    print(f"\n[{slide_num}] 创建结尾页...")
    create_ending_slide(prs, slide_num)

    # 保存
    print(f"\n保存PPT到: {OUTPUT_PATH}")
    prs.save(OUTPUT_PATH)
    print(f"\n✅ 完成! 共 {slide_num} 页幻灯片")
    print(f"   输出文件: {OUTPUT_PATH}")


if __name__ == '__main__':
    main()